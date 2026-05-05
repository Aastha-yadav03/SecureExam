from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Color Palette
DARK_BG   = RGBColor(0x0F, 0x17, 0x2A)
CARD_BG   = RGBColor(0x1E, 0x29, 0x3B)
PURPLE    = RGBColor(0x63, 0x66, 0xF1)
GREEN     = RGBColor(0x10, 0xB9, 0x81)
RED_ACCENT= RGBColor(0xF4, 0x3F, 0x5E)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
MUTED     = RGBColor(0x94, 0xA3, 0xB8)
LIGHT_BG  = RGBColor(0x1E, 0x29, 0x3B)

blank_layout = prs.slide_layouts[6]

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_card(slide, l, t, w, h, title, body, title_color=PURPLE):
    add_rect(slide, l, t, w, h, CARD_BG)
    add_rect(slide, l, t, 0.05, h, title_color)
    add_text(slide, title, l+0.15, t+0.15, w-0.3, 0.4, 13, bold=True, color=WHITE)
    add_text(slide, body, l+0.15, t+0.6, w-0.3, h-0.75, 11, color=MUTED)

# ── SLIDE 1: TITLE ──────────────────────────────────────────────
s1 = prs.slides.add_slide(blank_layout)
set_bg(s1, DARK_BG)
add_rect(s1, 0, 2.6, 13.33, 0.05, PURPLE)
add_rect(s1, 0, 0, 13.33, 7.5, RGBColor(0x0F,0x17,0x2A))
# Gradient accent bar
add_rect(s1, 0, 0, 0.6, 7.5, PURPLE)
add_rect(s1, 0, 0, 0.12, 7.5, RGBColor(0x81,0x8C,0xF8))
add_text(s1, "🔐 SecureExam", 1.0, 1.5, 11, 1.2, 52, bold=True, color=WHITE)
add_text(s1, "Exam Paper Security & Management System", 1.0, 2.9, 11, 0.8, 22, color=PURPLE)
add_text(s1, "Role-Based Access  •  AES-256 Encryption  •  Audit Trails  •  Real-Time Monitoring", 1.0, 3.8, 11, 0.6, 14, color=MUTED)
add_text(s1, "Built with React • Node.js • MongoDB", 1.0, 6.5, 8, 0.5, 12, color=MUTED, italic=True)

# ── SLIDE 2: PROBLEM STATEMENT ────────────────────────────────
s2 = prs.slides.add_slide(blank_layout)
set_bg(s2, DARK_BG)
add_rect(s2, 0, 0, 13.33, 1.2, CARD_BG)
add_text(s2, "The Problem", 0.5, 0.25, 12, 0.7, 30, bold=True, color=WHITE)
add_text(s2, "Why traditional exam management is a security risk", 0.5, 0.75, 12, 0.4, 14, color=MUTED)

problems = [
    ("📧", "Paper Leaks", "Exam papers shared via email or printed copies can be easily leaked before the exam date."),
    ("🔓", "No Access Control", "Anyone with access to shared drives or folders can view, edit, or copy exam content."),
    ("🕵️", "No Audit Trail", "There is no way to track who accessed, modified, or distributed exam papers."),
    ("📄", "No Encryption", "Sensitive academic content is stored in plain text, making it vulnerable to breaches."),
]
positions = [(0.4,1.4),(3.5,1.4),(6.6,1.4),(9.7,1.4)]
for (l,t),(icon,title,body) in zip(positions, problems):
    add_rect(s2, l, t, 3.0, 5.5, CARD_BG)
    add_rect(s2, l, t, 3.0, 0.05, RED_ACCENT)
    add_text(s2, icon, l+1.1, t+0.3, 0.8, 0.8, 30, align=PP_ALIGN.CENTER)
    add_text(s2, title, l+0.15, t+1.2, 2.7, 0.5, 14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s2, body, l+0.15, t+1.8, 2.7, 3.0, 11, color=MUTED, align=PP_ALIGN.CENTER)

# ── SLIDE 3: SOLUTION OVERVIEW ────────────────────────────────
s3 = prs.slides.add_slide(blank_layout)
set_bg(s3, DARK_BG)
add_rect(s3, 0, 0, 13.33, 1.2, CARD_BG)
add_text(s3, "Our Solution — SecureExam", 0.5, 0.25, 12, 0.7, 30, bold=True, color=WHITE)
add_text(s3, "A complete end-to-end secure platform for academic exam management", 0.5, 0.75, 12, 0.4, 14, color=MUTED)

add_rect(s3, 0.4, 1.4, 12.5, 5.5, CARD_BG)
add_text(s3, "SecureExam provides a centralized, encrypted, and auditable platform that allows faculty to create, encrypt, and publish exams — while giving administrators full visibility over who accessed what, and when.\n\nKey benefits:\n\n✅  AES-256 Encryption on all exam content\n✅  Role-based access: Admin, Faculty, Reviewer\n✅  Complete audit trail for every action\n✅  Multi-factor authentication (MFA) support\n✅  Real-time notifications via WebSocket\n✅  Digital watermarking of every exam paper\n✅  Suspicious activity detection & alerts",
         0.8, 1.6, 12.0, 5.0, 14, color=MUTED)
add_rect(s3, 0.4, 1.4, 0.07, 5.5, GREEN)

# ── SLIDE 4: SYSTEM ARCHITECTURE ─────────────────────────────
s4 = prs.slides.add_slide(blank_layout)
set_bg(s4, DARK_BG)
add_rect(s4, 0, 0, 13.33, 1.2, CARD_BG)
add_text(s4, "System Architecture", 0.5, 0.25, 12, 0.7, 30, bold=True, color=WHITE)
add_text(s4, "Three-tier full-stack architecture", 0.5, 0.75, 12, 0.4, 14, color=MUTED)

layers = [
    (PURPLE,    "🖥️  Frontend — React.js",   "• Login / Register Pages\n• Dashboard (Exam listing)\n• Exam Creator (Faculty)\n• Admin Panel\n• Context API for state management"),
    (GREEN,     "⚙️  Backend — Node.js / Express", "• REST API + WebSocket (Socket.IO)\n• JWT Authentication Middleware\n• Rate Limiting & Helmet Security\n• AES-256 Encryption Service\n• Audit Logging Service"),
    (RED_ACCENT,"🗄️  Database — MongoDB",    "• Users, ExamPapers, AccessLogs\n• AuditTrail, Alerts collections\n• In-memory DB for dev (no install needed)\n• Indexed queries for performance"),
]
for i,(col,title,body) in enumerate(layers):
    l = 0.4 + i*4.3
    add_rect(s4, l, 1.4, 4.1, 5.5, CARD_BG)
    add_rect(s4, l, 1.4, 4.1, 0.07, col)
    add_text(s4, title, l+0.15, 1.55, 3.8, 0.5, 13, bold=True, color=col)
    add_text(s4, body, l+0.15, 2.15, 3.8, 4.5, 12, color=MUTED)
    if i < 2:
        add_text(s4, "→", l+4.1, 3.8, 0.3, 0.5, 20, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

# ── SLIDE 5: SECURITY FEATURES ────────────────────────────────
s5 = prs.slides.add_slide(blank_layout)
set_bg(s5, DARK_BG)
add_rect(s5, 0, 0, 13.33, 1.2, CARD_BG)
add_text(s5, "Security Architecture", 0.5, 0.25, 12, 0.7, 30, bold=True, color=WHITE)
add_text(s5, "5 layers of protection built into every request", 0.5, 0.75, 12, 0.4, 14, color=MUTED)

sec_items = [
    (GREEN,     "Layer 1 — Input Protection",    "XSS sanitization, input validation, rate limiting (100 req/15 min)"),
    (PURPLE,    "Layer 2 — Authentication",       "bcryptjs password hashing, JWT tokens (7-day), MFA, account lockout after 5 attempts"),
    (RED_ACCENT,"Layer 3 — Encryption",           "AES-256-CBC for exam content, SHA-256 for integrity hashing"),
    (RGBColor(0xF5,0x9E,0x0B), "Layer 4 — Authorization", "Role-based access control (Admin/Faculty/Reviewer), endpoint protection"),
    (RGBColor(0x38,0xBD,0xF8), "Layer 5 — Logging",       "Access logs, audit trail, IP tracking, suspicious activity detection"),
]
for i,(col,title,body) in enumerate(sec_items):
    t = 1.4 + i*1.16
    add_rect(s5, 0.4, t, 12.5, 1.0, CARD_BG)
    add_rect(s5, 0.4, t, 0.08, 1.0, col)
    add_text(s5, title, 0.65, t+0.1, 4.0, 0.4, 13, bold=True, color=col)
    add_text(s5, body, 0.65, t+0.5, 12.0, 0.4, 11, color=MUTED)

# ── SLIDE 6: USER ROLES ───────────────────────────────────────
s6 = prs.slides.add_slide(blank_layout)
set_bg(s6, DARK_BG)
add_rect(s6, 0, 0, 13.33, 1.2, CARD_BG)
add_text(s6, "User Roles & Access Control", 0.5, 0.25, 12, 0.7, 30, bold=True, color=WHITE)
add_text(s6, "Three specialized roles, each with controlled permissions", 0.5, 0.75, 12, 0.4, 14, color=MUTED)

roles = [
    (RED_ACCENT, "👑 Administrator", "admin@secureexam.com", [
        "View all audit logs & access logs",
        "Monitor suspicious activities",
        "Create & manage user accounts",
        "View all exams across departments",
        "Generate security reports",
    ]),
    (PURPLE, "🎓 Faculty", "faculty@secureexam.com", [
        "Create new exam papers",
        "Add & encrypt questions",
        "Publish or archive exams",
        "View own exam history",
        "Digital watermark on each paper",
    ]),
    (GREEN, "🔍 Reviewer", "reviewer@secureexam.com", [
        "View published exam papers",
        "Review encrypted content",
        "View exam metadata",
        "Read-only access to exam list",
        "Cannot modify or create exams",
    ]),
]
for i,(col,role,email,perms) in enumerate(roles):
    l = 0.4 + i*4.3
    add_rect(s6, l, 1.4, 4.1, 5.5, CARD_BG)
    add_rect(s6, l, 1.4, 4.1, 0.07, col)
    add_text(s6, role, l+0.2, 1.55, 3.8, 0.5, 16, bold=True, color=col)
    add_text(s6, email, l+0.2, 2.1, 3.8, 0.35, 10, color=MUTED, italic=True)
    for j,p in enumerate(perms):
        add_text(s6, f"✓  {p}", l+0.2, 2.55+j*0.62, 3.8, 0.5, 11, color=WHITE)

# ── SLIDE 7: APPLICATION WORKFLOW ─────────────────────────────
s7 = prs.slides.add_slide(blank_layout)
set_bg(s7, DARK_BG)
add_rect(s7, 0, 0, 13.33, 1.2, CARD_BG)
add_text(s7, "Application Workflow", 0.5, 0.25, 12, 0.7, 30, bold=True, color=WHITE)
add_text(s7, "From login to secure exam creation — step by step", 0.5, 0.75, 12, 0.4, 14, color=MUTED)

steps = [
    (PURPLE,    "1", "Login",        "Enter email & password on the secure login page"),
    (GREEN,     "2", "Authenticate", "Backend verifies credentials, issues JWT token"),
    (RED_ACCENT,"3", "Dashboard",    "User sees role-specific dashboard with their exams"),
    (RGBColor(0xF5,0x9E,0x0B),"4","Create Exam","Faculty fills exam details, adds questions"),
    (RGBColor(0x38,0xBD,0xF8),"5","Encrypt & Save","Backend encrypts with AES-256, generates hash & watermark"),
    (GREEN,     "6", "Publish",      "Exam is published; audit trail entry is created automatically"),
]
for i,(col,num,title,body) in enumerate(steps):
    col_pos = i % 3
    row_pos = i // 3
    l = 0.4 + col_pos*4.3
    t = 1.4 + row_pos*2.8
    add_rect(s7, l, t, 4.1, 2.5, CARD_BG)
    add_rect(s7, l, t, 0.07, 2.5, col)
    add_rect(s7, l+0.2, t+0.2, 0.5, 0.5, col)
    add_text(s7, num, l+0.2, t+0.2, 0.5, 0.5, 18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s7, title, l+0.85, t+0.25, 3.0, 0.4, 14, bold=True, color=col)
    add_text(s7, body, l+0.2, t+0.85, 3.7, 1.4, 11, color=MUTED)

# ── SLIDE 8: KEY FEATURES ─────────────────────────────────────
s8 = prs.slides.add_slide(blank_layout)
set_bg(s8, DARK_BG)
add_rect(s8, 0, 0, 13.33, 1.2, CARD_BG)
add_text(s8, "Key Features", 0.5, 0.25, 12, 0.7, 30, bold=True, color=WHITE)
add_text(s8, "Everything you need for secure academic assessment", 0.5, 0.75, 12, 0.4, 14, color=MUTED)

features = [
    ("🔐","AES-256 Encryption","All exam content is encrypted before storage"),
    ("🪪","JWT Authentication","Secure token-based session management"),
    ("📋","Full Audit Trail","Every action is logged with timestamp & IP"),
    ("📡","Real-time Updates","WebSocket notifications for live collaboration"),
    ("🛡️","MFA Support","Optional Multi-Factor Authentication per user"),
    ("🔏","Digital Watermark","userId + timestamp embedded in every paper"),
    ("⚠️","Anomaly Detection","Suspicious activity flagged automatically"),
    ("👥","Role-Based Access","Admin, Faculty, Reviewer — granular permissions"),
    ("📊","Admin Dashboard","Full visibility over logs and user activity"),
]
for i,(icon,title,body) in enumerate(features):
    col_pos = i % 3
    row_pos = i // 3
    l = 0.4 + col_pos*4.3
    t = 1.4 + row_pos*1.9
    add_rect(s8, l, t, 4.1, 1.7, CARD_BG)
    add_text(s8, icon, l+0.15, t+0.2, 0.6, 0.6, 18)
    add_text(s8, title, l+0.8, t+0.2, 3.1, 0.45, 13, bold=True, color=WHITE)
    add_text(s8, body, l+0.8, t+0.7, 3.1, 0.8, 11, color=MUTED)

# ── SLIDE 9: TECH STACK ───────────────────────────────────────
s9 = prs.slides.add_slide(blank_layout)
set_bg(s9, DARK_BG)
add_rect(s9, 0, 0, 13.33, 1.2, CARD_BG)
add_text(s9, "Technology Stack", 0.5, 0.25, 12, 0.7, 30, bold=True, color=WHITE)
add_text(s9, "Modern, battle-tested technologies powering SecureExam", 0.5, 0.75, 12, 0.4, 14, color=MUTED)

tech = [
    (PURPLE,    "Frontend", [("React.js","SPA framework"),("React Router","Client-side routing"),("Context API","State management"),("Socket.IO Client","Real-time updates"),("Outfit Font","Premium typography")]),
    (GREEN,     "Backend",  [("Node.js","Server runtime"),("Express.js","REST API framework"),("Socket.IO","WebSocket server"),("JWT","Auth tokens"),("bcryptjs","Password hashing")]),
    (RED_ACCENT,"Security", [("AES-256-CBC","Exam encryption"),("SHA-256","Integrity hashing"),("Helmet.js","HTTP security headers"),("Rate Limiter","DDoS protection"),("CORS","Origin control")]),
    (RGBColor(0xF5,0x9E,0x0B),"Database", [("MongoDB","Primary database"),("Mongoose","ODM layer"),("In-Memory DB","Dev mode (no install)"),("Indexing","Performance queries"),("Audit Collections","Log storage")]),
]
for i,(col,cat,items) in enumerate(tech):
    l = 0.4 + i*3.2
    add_rect(s9, l, 1.4, 3.0, 5.6, CARD_BG)
    add_rect(s9, l, 1.4, 3.0, 0.07, col)
    add_text(s9, cat, l+0.15, 1.55, 2.7, 0.45, 14, bold=True, color=col)
    for j,(name,desc) in enumerate(items):
        t = 2.1 + j*0.92
        add_rect(s9, l+0.15, t, 2.7, 0.8, RGBColor(0x0F,0x17,0x2A))
        add_text(s9, name, l+0.3, t+0.05, 2.4, 0.35, 12, bold=True, color=WHITE)
        add_text(s9, desc, l+0.3, t+0.42, 2.4, 0.3, 10, color=MUTED)

# ── SLIDE 10: THANK YOU ───────────────────────────────────────
s10 = prs.slides.add_slide(blank_layout)
set_bg(s10, DARK_BG)
add_rect(s10, 0, 0, 0.6, 7.5, PURPLE)
add_rect(s10, 0, 0, 0.12, 7.5, RGBColor(0x81,0x8C,0xF8))
add_rect(s10, 0.6, 3.3, 12.73, 0.05, PURPLE)
add_text(s10, "Thank You!", 1.0, 1.0, 11, 1.5, 60, bold=True, color=WHITE)
add_text(s10, "SecureExam — Securing Academic Integrity", 1.0, 2.7, 11, 0.7, 22, color=PURPLE)
add_text(s10, "🔐  Encrypted  •  🛡️  Secure  •  📊  Auditable  •  👥  Role-Based", 1.0, 3.6, 11, 0.6, 16, color=MUTED)

add_card(s10, 1.0, 4.4, 3.5, 2.5, "Default Login Credentials", "Admin:\nadmin@secureexam.com\nAdminPass123!\n\nFaculty:\nfaculty@secureexam.com\nFacultyPass123!", title_color=RED_ACCENT)
add_card(s10, 4.8, 4.4, 3.5, 2.5, "Access the App", "Frontend:\nhttp://localhost:3000\n\nBackend API:\nhttp://localhost:5000/api", title_color=GREEN)
add_card(s10, 8.6, 4.4, 3.9, 2.5, "Tech Stack", "React • Node.js • Express\nMongoDB • Socket.IO\nJWT • AES-256 • bcryptjs", title_color=PURPLE)

# Save
output = r"c:\Users\91921.LAPTOP-30HKO39J\Desktop\SecureExam_Presentation.pptx"
prs.save(output)
print(f"Presentation saved to: {output}")
